"""
FastAPI backend for RAG chatbot UI
"""
from fastapi import FastAPI, HTTPException, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import sqlite3
import pickle
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from rank_bm25 import BM25Okapi
import subprocess
import datetime
import textwrap
import os
import shutil
from pathlib import Path

# -----------------------
# Configuration
# -----------------------
DB_FILE = "chunks.db"
EMBEDDING_MODEL = "all-MiniLM-L6-v2"
TOP_K_FAISS = 100
TOP_K_BM25 = 100
FINAL_TOP_K = 8
DEFAULT_ALPHA = 0.6
MIN_COMBINED_SCORE_TO_ANSWER = 0.12
OLLAMA_MODEL = "gpt-oss:20b-cloud"
ADMIN_PASSWORD = "admin@123"
UPLOAD_DIRECTORY = "uploads"
CHUNKING_SCRIPT = "scripts/chunking_script.py"

# Create upload directory if it doesn't exist
Path(UPLOAD_DIRECTORY).mkdir(exist_ok=True)

# -----------------------
# Models
# -----------------------
class QuestionRequest(BaseModel):
    session_id: Optional[int] = None
    question: str
    use_cache: bool = True

class SessionCreate(BaseModel):
    name: str = "New Chat"

class SessionInfo(BaseModel):
    session_id: int
    name: str
    created_at: str
    message_count: int

class Message(BaseModel):
    role: str
    message: str
    timestamp: str

class ChatResponse(BaseModel):
    answer: str
    session_id: int
    sources: List[Dict[str, Any]]

class AdminUploadRequest(BaseModel):
    password: str

# -----------------------
# RAG System Class
# -----------------------
class RAGSystemAPI:
    def __init__(self, db_path=DB_FILE, embedding_model=EMBEDDING_MODEL):
        if not os.path.exists(db_path):
            raise FileNotFoundError(f"DB file '{db_path}' not found.")
        
        self.conn = sqlite3.connect(db_path, check_same_thread=False)
        self.init_management_tables()
        
        print("Loading chunks and building indices...")
        self.chunks = self.load_chunks()
        print(f"Loaded {len(self.chunks)} chunks.")
        
        self.faiss_index, self.emb_matrix = self.build_faiss_from_chunks(self.chunks)
        self.bm25, self.bm25_tokenized = self.build_bm25_from_chunks(self.chunks)
        
        self.embed_model = SentenceTransformer(embedding_model)
        self.alpha = DEFAULT_ALPHA
    
    def init_management_tables(self):
        cur = self.conn.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chat_sessions (
                session_id INTEGER PRIMARY KEY AUTOINCREMENT,
                name TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chat_messages (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id INTEGER,
                role TEXT,
                message TEXT,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY(session_id) REFERENCES chat_sessions(session_id)
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS question_cache (
                question TEXT PRIMARY KEY,
                answer TEXT,
                saved_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        self.conn.commit()
    
    def load_chunks(self):
        chunks = []
        cur = self.conn.cursor()
        cur.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [r[0] for r in cur.fetchall()]
        skip = {"sqlite_sequence", "chat_sessions", "chat_messages", "question_cache"}
        tables = [t for t in tables if t not in skip]
        
        for t in tables:
            try:
                cur.execute(f'SELECT chunk_id, page_number, chunk_text, embedding FROM "{t}"')
                rows = cur.fetchall()
                for row in rows:
                    cid, pnum, ctext, emb_blob = row
                    if emb_blob is None:
                        continue
                    try:
                        emb = pickle.loads(emb_blob)
                    except:
                        continue
                    chunks.append({
                        "table": t,
                        "chunk_id": int(cid),
                        "page_number": int(pnum) if pnum is not None else None,
                        "text": ctext,
                        "embedding": np.array(emb, dtype=np.float32)
                    })
            except:
                continue
        return chunks
    
    def build_faiss_from_chunks(self, chunks):
        if len(chunks) == 0:
            d = 384
            index = faiss.IndexFlatIP(d)
            return index, np.zeros((0, d), dtype=np.float32)
        embs = np.vstack([c["embedding"].astype(np.float32) for c in chunks])
        faiss.normalize_L2(embs)
        index = faiss.IndexFlatIP(embs.shape[1])
        index.add(embs)
        return index, embs
    
    def build_bm25_from_chunks(self, chunks):
        corpus = [c["text"] if c["text"] else "" for c in chunks]
        tokenized = [doc.split() for doc in corpus]
        bm25 = BM25Okapi(tokenized)
        return bm25, tokenized
    
    def normalize_scores(self, arr):
        if arr.size == 0:
            return arr
        mn = float(np.min(arr))
        mx = float(np.max(arr))
        if abs(mx - mn) < 1e-12:
            return np.zeros_like(arr)
        return (arr - mn) / (mx - mn)
    
    def hybrid_retrieve(self, query, top_k_faiss=TOP_K_FAISS, top_k_bm25=TOP_K_BM25, final_k=FINAL_TOP_K):
        q_emb = self.embed_model.encode([query], convert_to_numpy=True).astype(np.float32)
        faiss.normalize_L2(q_emb)
        
        if self.faiss_index.ntotal == 0:
            faiss_scores = np.zeros(len(self.chunks), dtype=np.float32)
        else:
            k = min(top_k_faiss, self.faiss_index.ntotal)
            scores, ids = self.faiss_index.search(q_emb, k)
            faiss_scores = np.zeros(len(self.chunks), dtype=np.float32)
            if ids is not None and scores is not None:
                ids = ids.flatten()
                scores = scores.flatten()
                faiss_scores[ids] = scores
        
        tokenized_q = query.split()
        bm25_scores = np.array(self.bm25.get_scores(tokenized_q), dtype=np.float32)
        
        f_norm = self.normalize_scores(faiss_scores)
        b_norm = self.normalize_scores(bm25_scores)
        
        combined = self.alpha * f_norm + (1.0 - self.alpha) * b_norm
        
        top_idxs = np.argsort(-combined)[:final_k]
        results = []
        for idx in top_idxs:
            if combined[idx] <= 0:
                continue
            chunk_info = {
                "table": self.chunks[idx]["table"],
                "chunk_id": self.chunks[idx]["chunk_id"],
                "page": self.chunks[idx]["page_number"],
                "text": self.chunks[idx]["text"][:300] + "...",
                "score": float(combined[idx]),
                "faiss_score": float(f_norm[idx]),
                "bm25_score": float(b_norm[idx])
            }
            results.append(chunk_info)
        
        return sorted(results, key=lambda x: x["score"], reverse=True)
    
    def call_ollama(self, prompt):
        try:
            proc = subprocess.run(
                ["ollama", "run", OLLAMA_MODEL],
                input=prompt.encode("utf-8"),
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=True,
                timeout=300
            )
            out = proc.stdout.decode("utf-8").strip()
            if not out:
                err = proc.stderr.decode("utf-8").strip()
                return err
            return out
        except Exception as e:
            raise RuntimeError(f"Ollama call error: {e}")
    
    def create_session(self, name="New Chat"):
        cur = self.conn.cursor()
        cur.execute("INSERT INTO chat_sessions (name) VALUES (?)", (name,))
        self.conn.commit()
        return cur.lastrowid
    
    def list_sessions(self):
        cur = self.conn.cursor()
        cur.execute("""
            SELECT s.session_id, s.name, s.created_at, 
                   COUNT(m.id) as message_count
            FROM chat_sessions s
            LEFT JOIN chat_messages m ON s.session_id = m.session_id
            GROUP BY s.session_id
            ORDER BY s.created_at DESC
        """)
        return cur.fetchall()
    
    def get_session_messages(self, session_id):
        cur = self.conn.cursor()
        cur.execute("""
            SELECT role, message, timestamp 
            FROM chat_messages 
            WHERE session_id = ? 
            ORDER BY id ASC
        """, (session_id,))
        return cur.fetchall()
    
    def delete_session(self, session_id):
        cur = self.conn.cursor()
        cur.execute("DELETE FROM chat_messages WHERE session_id = ?", (session_id,))
        cur.execute("DELETE FROM chat_sessions WHERE session_id = ?", (session_id,))
        self.conn.commit()
    
    def append_message(self, session_id, role, message):
        cur = self.conn.cursor()
        cur.execute("""
            INSERT INTO chat_messages (session_id, role, message) 
            VALUES (?, ?, ?)
        """, (session_id, role, message))
        self.conn.commit()
    
    def ask_question(self, session_id, question, use_cache=True):
        # Retrieve relevant chunks
        retrieved = self.hybrid_retrieve(question)
        
        if not retrieved or retrieved[0]["score"] < MIN_COMBINED_SCORE_TO_ANSWER:
            response = "Data not found"
            self.append_message(session_id, "user", question)
            self.append_message(session_id, "assistant", response)
            return {
                "answer": response,
                "sources": [],
                "from_cache": False
            }
        
        # Build prompt for detailed synthesis with table support
        history = self.get_session_messages(session_id)
        history_text = ""
        if history:
            history_text = "Previous conversation (for context only):\n"
            for role, msg, ts in history[-3:]:
                history_text += f"{role}: {msg}\n"
        
        # Prepare context chunks with clear source information
        chunks_with_sources = []
        for i, r in enumerate(retrieved[:5]):
            source_info = f"[Source {i+1}: {r['table']}, Page {r['page']}]"
            chunk_content = r['text'].replace('\n', ' ').strip()
            chunks_with_sources.append(f"{source_info}: {chunk_content}")
        
        context_text = "\n\n".join(chunks_with_sources)
        
        # Enhanced prompt for detailed synthesis with markdown and table support
        prompt = f"""# DOCUMENT-BASED QUESTION ANSWERING TASK

## CONTEXT FROM DOCUMENTS:
{context_text}

## QUESTION TO ANSWER:
{question}

## INSTRUCTIONS:
1. **USE ONLY THE PROVIDED CONTEXT ABOVE** - Do not use any external knowledge
2. **SYNTHESIZE INFORMATION** - Combine relevant information from multiple sources if applicable
3. **BE DETAILED AND WELL-STRUCTURED** - Provide comprehensive answer with clear organization
4. **USE MARKDOWN FORMATTING**:
   - Use **bold** for emphasis and important terms
   - Use bullet points for lists
   - Use tables for comparisons, data presentation, or structured information
   - Use headers (##, ###) for section organization
5. **TABLE CREATION GUIDELINES**:
   - Create tables when comparing multiple items, showing features, or presenting structured data
   - Use proper markdown table syntax with headers and alignment
   - Example table format:
     | Feature | Item A | Item B | Item C |
     |---------|--------|--------|--------|
     | Price   | $100   | $150   | $200   |
     | Rating  | 4.5    | 4.2    | 4.8    |
6. **CITE YOUR SOURCES** - For each key point, include citation like [Source X]
7. **IF INFORMATION IS INSUFFICIENT** - If the context does not contain enough information to answer the question properly, respond with exactly: "Data not found"
8. **DO NOT INVENT INFORMATION** - Only use what's explicitly stated in the context
9. **DO NOT INCLUDE THINKING OR ANALYSIS** - Provide only the final answer

## ANSWER STRUCTURE:
1. Start with a direct, concise answer to the question
2. Provide detailed explanations with supporting evidence from sources
3. Use tables when presenting comparisons, features, or structured data
4. Include specific examples from the context with proper citations
5. End with a brief summary or conclusion
6. Format using markdown for better readability

## SYNTHESIZED ANSWER (with markdown formatting):
Based ONLY on the provided context, provide a comprehensive answer to the question. Use markdown formatting including tables where appropriate:"""
        
        try:
            answer = self.call_ollama(prompt)
            
            # Post-process the answer
            import re
            answer = answer.strip()
            
            print(f"DEBUG: Raw answer from Ollama:\n{answer[:500]}...\n")
            
            # Remove thinking/analysis XML tags if present
            answer = re.sub(r'<thinking>.*?</thinking>', '', answer, flags=re.DOTALL | re.IGNORECASE)
            answer = re.sub(r'\[thinking\].*?\[/thinking\]', '', answer, flags=re.DOTALL | re.IGNORECASE)
            
            print(f"DEBUG: After removing XML tags:\n{answer[:500]}...\n")
            
            # Remove thinking preamble lines from the start
            lines = answer.split('\n')
            trimmed_lines = []
            skip = True
            
            for line in lines:
                line_strip = line.strip()
                line_lower = line_strip.lower()
                
                # Skip thinking/analysis preamble at the beginning
                if skip and line_strip:
                    if any(line_lower.startswith(prefix) for prefix in 
                           ['thinking:', 'analysis:', 'let me think', 'let me analyze', 'hmm', 'well,']):
                        print(f"DEBUG: Skipping preamble line: {line_strip[:100]}")
                        continue
                    skip = False
                
                # Keep non-empty lines after preamble
                trimmed_lines.append(line)
            
            answer = '\n'.join(trimmed_lines).strip()
            
            print(f"DEBUG: After removing preamble:\n{answer[:500]}...\n")
            
            # Ensure citations are present and clean
            if answer and answer != "Data not found":
                # Check if we have any citations
                has_citations = any(f"[Source" in answer for f in ["[Source", "[Table:"])
                if not has_citations and retrieved:
                    top_source = retrieved[0]
                    answer = f"{answer}\n\n*Based on information from: {top_source['table']}, Page {top_source['page']}*"
            
            # Ensure we don't have empty answers
            if not answer or answer.isspace():
                answer = "Data not found"
                
        except Exception as e:
            answer = f"Error generating answer: {str(e)}"
        
        # Save to chat history
        self.append_message(session_id, "user", question)
        self.append_message(session_id, "assistant", answer)
        
        return {
            "answer": answer,
            "sources": retrieved[:5],
            "from_cache": False
        }

# -----------------------
# FastAPI App
# -----------------------
app = FastAPI(title="RAG Chatbot API", version="1.0.0")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:3001"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global RAG instance
rag_system = None

@app.on_event("startup")
async def startup_event():
    global rag_system
    try:
        rag_system = RAGSystemAPI()
        print("RAG system initialized successfully")
    except Exception as e:
        print(f"Failed to initialize RAG system: {e}")
        raise

# -----------------------
# API Endpoints
# -----------------------
@app.get("/")
async def root():
    return {"message": "RAG Chatbot API", "status": "running"}

@app.get("/api/health")
async def health_check():
    return {"status": "healthy", "chunks_loaded": len(rag_system.chunks)}

@app.post("/api/sessions", response_model=Dict[str, Any])
async def create_session(session: SessionCreate):
    session_id = rag_system.create_session(session.name)
    return {"session_id": session_id, "name": session.name, "status": "created"}

@app.get("/api/sessions", response_model=List[SessionInfo])
async def get_sessions():
    sessions = rag_system.list_sessions()
    return [
        SessionInfo(
            session_id=row[0],
            name=row[1],
            created_at=row[2],
            message_count=row[3]
        )
        for row in sessions
    ]

@app.get("/api/sessions/{session_id}/messages", response_model=List[Message])
async def get_messages(session_id: int):
    messages = rag_system.get_session_messages(session_id)
    # Ensure timestamps are properly formatted
    formatted_messages = []
    for role, message, timestamp in messages:
        # Convert timestamp to ISO format if it's not already
        if timestamp and not isinstance(timestamp, str):
            timestamp = str(timestamp)
        formatted_messages.append({
            "role": role,
            "message": message,
            "timestamp": timestamp or datetime.datetime.utcnow().isoformat()
        })
    return formatted_messages

@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: int):
    rag_system.delete_session(session_id)
    return {"status": "deleted", "session_id": session_id}

@app.post("/api/ask", response_model=ChatResponse)
async def ask_question(request: QuestionRequest):
    if not request.session_id:
        request.session_id = rag_system.create_session("Auto-created session")
    
    result = rag_system.ask_question(
        request.session_id,
        request.question,
        request.use_cache
    )
    
    return ChatResponse(
        answer=result["answer"],
        session_id=request.session_id,
        sources=result["sources"]
    )

@app.put("/api/alpha/{alpha_value}")
async def set_alpha(alpha_value: float):
    if 0.0 <= alpha_value <= 1.0:
        rag_system.alpha = alpha_value
        return {"alpha": rag_system.alpha, "status": "updated"}
    raise HTTPException(status_code=400, detail="Alpha must be between 0 and 1")

@app.get("/api/stats")
async def get_stats():
    return {
        "total_chunks": len(rag_system.chunks),
        "alpha": rag_system.alpha,
        "tables": len(set(c["table"] for c in rag_system.chunks))
    }

@app.post("/api/admin/upload")
async def admin_upload(password: str = Form(...), files: List[UploadFile] = File(...)):
    """
    Admin endpoint for uploading documents to be chunked.
    Password protected with admin@123
    Supported formats: PDF, DOCX, PPTX
    """
    # Verify admin password
    if password != ADMIN_PASSWORD:
        raise HTTPException(status_code=403, detail="Invalid admin password")
    
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    
    try:
        # Save uploaded files to temp directory
        saved_files = []
        for file in files:
            if not file.filename:
                continue
                
            # Check file extension
            file_ext = Path(file.filename).suffix.lower()
            if file_ext not in ['.pdf', '.docx', '.pptx']:
                raise HTTPException(
                    status_code=400, 
                    detail=f"Unsupported file type: {file_ext}. Supported: PDF, DOCX, PPTX"
                )
            
            # Save file
            file_path = Path(UPLOAD_DIRECTORY) / file.filename
            with open(file_path, "wb") as buffer:
                content = await file.read()
                buffer.write(content)
            saved_files.append(str(file_path))
        
        if not saved_files:
            raise HTTPException(status_code=400, detail="No valid files to process")
        
        # Execute chunking script
        try:
            result = execute_chunking(saved_files)
            
            # Clean up uploaded files after processing
            for file_path in saved_files:
                try:
                    os.remove(file_path)
                except:
                    pass
            
            # Reload the RAG system to include new chunks
            global rag_system
            rag_system.chunks = rag_system.load_chunks()
            rag_system.faiss_index, rag_system.emb_matrix = rag_system.build_faiss_from_chunks(rag_system.chunks)
            rag_system.bm25, rag_system.bm25_tokenized = rag_system.build_bm25_from_chunks(rag_system.chunks)
            
            return {
                "status": "success",
                "message": f"Processed {len(saved_files)} file(s)",
                "files": [Path(f).name for f in saved_files],
                "chunking_output": result
            }
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Chunking error: {str(e)}")
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Upload error: {str(e)}")

@app.get("/api/admin/documents")
async def get_documents(password: str):
    """
    Get list of uploaded documents (tables in the database)
    """
    if password != ADMIN_PASSWORD:
        raise HTTPException(status_code=403, detail="Invalid admin password")
    
    try:
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [r[0] for r in cursor.fetchall()]
        conn.close()
        
        # Filter out system tables
        skip = {"sqlite_sequence", "chat_sessions", "chat_messages", "question_cache"}
        document_tables = [t for t in tables if t not in skip]
        
        # Get chunk count for each document
        documents = []
        for table_name in document_tables:
            try:
                conn = sqlite3.connect(DB_FILE)
                cursor = conn.cursor()
                cursor.execute(f'SELECT COUNT(*) FROM "{table_name}"')
                chunk_count = cursor.fetchone()[0]
                conn.close()
                documents.append({
                    "name": table_name,
                    "chunks": chunk_count
                })
            except:
                pass
        
        return {"documents": documents}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching documents: {str(e)}")

@app.delete("/api/admin/documents/{doc_name}")
async def delete_document(doc_name: str, password: str):
    """
    Delete a document from the database
    """
    if password != ADMIN_PASSWORD:
        raise HTTPException(status_code=403, detail="Invalid admin password")
    
    try:
        # Validate table name to prevent SQL injection
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [r[0] for r in cursor.fetchall()]
        conn.close()
        
        if doc_name not in tables:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Drop the document table
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        cursor.execute(f'DROP TABLE "{doc_name}"')
        conn.commit()
        conn.close()
        
        # Reload the RAG system
        global rag_system
        rag_system.chunks = rag_system.load_chunks()
        rag_system.faiss_index, rag_system.emb_matrix = rag_system.build_faiss_from_chunks(rag_system.chunks)
        rag_system.bm25, rag_system.bm25_tokenized = rag_system.build_bm25_from_chunks(rag_system.chunks)
        
        return {"status": "deleted", "document": doc_name}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting document: {str(e)}")

def execute_chunking(file_paths: List[str]) -> str:
    """
    Execute the chunking script on uploaded files
    """
    # Create temporary config for chunking script
    directory_arg = " ".join(file_paths)
    
    # Since our files are already in the right place, we'll modify the 
    # chunking script call to use the uploaded files directly
    import sys
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    
    # Import and run chunking logic
    try:
        # Execute chunking for each file
        output_messages = []
        for file_path in file_paths:
            file_path = Path(file_path)
            if file_path.exists():
                output_messages.append(f"Processing: {file_path.name}")
        
        # Call the chunking script as a subprocess
        cmd = [
            "python", 
            CHUNKING_SCRIPT,
            "--files", 
            " ".join(file_paths)
        ]
        
        # For now, we'll execute the chunking inline to ensure it works with the DB
        # This avoids subprocess complexity with file paths
        exec_chunking_inline(file_paths)
        
        return f"Successfully processed {len(file_paths)} file(s)"
    except Exception as e:
        raise Exception(f"Chunking execution failed: {str(e)}")

def exec_chunking_inline(file_paths: List[str]):
    """
    Execute chunking logic directly without subprocess
    This ensures proper database integration
    """
    import nltk
    from sentence_transformers import SentenceTransformer
    from sklearn.metrics.pairwise import cosine_similarity
    import fitz  # PyMuPDF for PDFs
    from docx import Document
    from pptx import Presentation
    from tqdm import tqdm
    
    try:
        nltk.download('punkt', quiet=True)
    except:
        pass
    
    # Initialize embedding model
    model = SentenceTransformer(EMBEDDING_MODEL)
    
    # Connect to database
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    
    SIMILARITY_THRESHOLD = 0.9
    
    def serialize_embedding(embedding: np.ndarray) -> bytes:
        return pickle.dumps(embedding.astype(np.float32))
    
    def chunk_text_by_similarity(text, threshold=SIMILARITY_THRESHOLD):
        sentences = nltk.sent_tokenize(text)
        if len(sentences) == 0:
            return []
        
        embeddings = model.encode(sentences, convert_to_numpy=True, show_progress_bar=False)
        chunks = []
        current_chunk = [sentences[0]]
        current_embeds = [embeddings[0]]
        
        for i in range(1, len(sentences)):
            sim = cosine_similarity([embeddings[i - 1]], [embeddings[i]])[0][0]
            if sim >= threshold:
                current_chunk.append(sentences[i])
                current_embeds.append(embeddings[i])
            else:
                chunk_text = " ".join(current_chunk)
                chunk_embed = np.mean(current_embeds, axis=0)
                chunks.append((chunk_text, chunk_embed))
                current_chunk = [sentences[i]]
                current_embeds = [embeddings[i]]
        
        if current_chunk:
            chunk_text = " ".join(current_chunk)
            chunk_embed = np.mean(current_embeds, axis=0)
            chunks.append((chunk_text, chunk_embed))
        
        return chunks
    
    def extract_text_from_pdf(filepath):
        doc = fitz.open(filepath)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            pages.append((i + 1, text.strip()))
        doc.close()
        return pages
    
    def extract_text_from_docx(filepath):
        doc = Document(filepath)
        return [(i + 1, p.text.strip()) for i, p in enumerate(doc.paragraphs) if p.text.strip()]
    
    def extract_text_from_pptx(filepath):
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            slides.append((i + 1, "\n".join(text_runs).strip()))
        return slides
    
    def create_table(table_name):
        cursor.execute(f"""
            CREATE TABLE IF NOT EXISTS "{table_name}" (
                chunk_id INTEGER PRIMARY KEY AUTOINCREMENT,
                page_number INTEGER,
                chunk_text TEXT,
                embedding BLOB
            )
        """)
        conn.commit()
    
    # Process each file
    for file_path in file_paths:
        file_path = Path(file_path)
        if not file_path.exists():
            continue
        
        filename = file_path.name
        table_name = file_path.stem.replace(" ", "_").replace("-", "_")
        
        print(f"Processing: {filename}")
        create_table(table_name)
        
        # Extract text based on file type
        if str(file_path).lower().endswith(".pdf"):
            pages = extract_text_from_pdf(str(file_path))
        elif str(file_path).lower().endswith(".docx"):
            pages = extract_text_from_docx(str(file_path))
        elif str(file_path).lower().endswith(".pptx"):
            pages = extract_text_from_pptx(str(file_path))
        else:
            continue
        
        # Chunk each page
        for page_num, page_text in pages:
            if not page_text:
                continue
            
            chunks = chunk_text_by_similarity(page_text)
            
            for chunk_text, chunk_embed in chunks:
                blob = serialize_embedding(chunk_embed)
                cursor.execute(
                    f'INSERT INTO "{table_name}" (page_number, chunk_text, embedding) VALUES (?, ?, ?)',
                    (page_num, chunk_text, blob)
                )
        conn.commit()
    
    conn.close()
    print("✅ All files processed successfully!")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)